from flask import Flask, render_template, jsonify, send_from_directory, request, redirect, url_for, session, flash, \
    abort
from config import Config
from utils.auth import auth_manager, login_required, rate_limit
import bcrypt
import logging
from datetime import datetime
from functools import wraps
import struct
import io
import tempfile
import striprtf


app = Flask(__name__)
app.config.from_object(Config)

# Настройка логирования
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Хранилище для хешированных паролей
user_hashes = {}
user_roles = {}
user_names = {}


def init_users():
    """Инициализация пользователей с хешированием паролей"""
    for username, password in Config.USERS.items():
        salt = bcrypt.gensalt()
        user_hashes[username] = bcrypt.hashpw(password.encode('utf-8'), salt)
        # Определяем роль на основе имени пользователя
        if username == 'admin':
            user_roles[username] = 'admin'
            user_names[username] = 'Администратор'
        elif username == 'manager':
            user_roles[username] = 'manager'
            user_names[username] = 'Менеджер'
        elif username == 'guest':
            user_roles[username] = 'guest'
            user_names[username] = 'Гость'
        else:
            user_roles[username] = 'user'
            user_names[username] = username.capitalize()
    logger.info(f'Initialized {len(user_hashes)} users')


# Инициализация пользователей при старте
init_users()

# Основные данные
coordinates = [
    {"id": 1, "uid": 12, "latitude": 55.7558, "longitude": 37.6173, "name": "Москва"},
    {"id": 2, "uid": 13, "latitude": 59.9343, "longitude": 30.3351, "name": "Санкт-Петербург"},
    {"id": 3, "uid": 12, "latitude": 48.8566, "longitude": 2.3522, "name": "Париж"},
    {"id": 4, "uid": 17, "latitude": 51.5074, "longitude": -0.1278, "name": "Лондон"},
    {"id": 5, "uid": 12, "latitude": -40.7128, "longitude": -74.0060, "name": "Нью-Йорк"},
    {"id": 6, "uid": 122, "latitude": 35.6762, "longitude": 139.6503, "name": "Токио"},
    {"id": 7, "uid": 12, "latitude": -33.8688, "longitude": 151.2093, "name": "Сидней"},
    {"id": 8, "uid": 11, "latitude": -23.5505, "longitude": -46.6333, "name": "Сан-Паулу"},
    {"id": 9, "uid": 32, "latitude": 1.0760, "longitude": 72.8777, "name": "Мумбаи"},
    {"id": 10, "uid": 2, "latitude": 0.0444, "longitude": 71.2357, "name": "Каир"},
    {"id": 11, "uid": 122, "latitude": 50.04, "longitude": 30.2357, "name": "Киев"},
    {"id": 12, "uid": 112, "latitude": 50.04, "longitude": 33.2357, "name": "Киев"},
    {"id": 13, "uid": 12, "latitude": 50.04, "longitude": 31.2357, "name": "Киев"},
    {"id": 14, "uid": 125, "latitude": 30.04, "longitude": 30.2357, "name": "Киев"},
    {"id": 15, "uid": 126, "latitude": 55.04, "longitude": 35.2357, "name": "Киев"},
    {"id": 16, "uid": 1, "latitude": 10.04, "longitude": 38.2357, "name": "Киев"},
]


# ============================================
# ДЕКОРАТОРЫ ДЛЯ ПРОВЕРКИ РОЛЕЙ
# ============================================

def role_required(allowed_roles):
    """
    Декоратор для проверки роли пользователя
    Использование: @role_required(['admin', 'manager'])
    """

    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            # Проверяем авторизацию
            if not session.get('logged_in'):
                abort(401)

            # Получаем роль пользователя
            user_role = session.get('user_role', 'guest')

            # Проверяем доступ
            if user_role not in allowed_roles:
                flash('У вас недостаточно прав для доступа к этой странице.', 'danger')
                abort(403)

            return f(*args, **kwargs)

        return decorated_function

    return decorator


def role_required_from_config():
    """
    Декоратор для проверки роли на основе конфигурации
    Использование: @role_required_from_config()
    """

    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            if not session.get('logged_in'):
                abort(401)

            # Получаем текущий URL
            url = request.path

            # Проверяем права в конфигурации
            allowed_roles = Config.PAGE_ACCESS.get(url, ['admin', 'user', 'manager', 'guest'])
            user_role = session.get('user_role', 'guest')

            if user_role not in allowed_roles:
                flash('У вас недостаточно прав для доступа к этой странице.', 'danger')
                abort(403)

            return f(*args, **kwargs)

        return decorated_function

    return decorator


# ============================================
# МАРШРУТЫ
# ============================================

@app.route('/')
@login_required
def index():
    """Главная страница - доступна всем авторизованным"""
    return render_template('index.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


@app.route('/dashboard')
@login_required
@role_required_from_config()
def dashboard():
    """Дашборд - для всех кроме гостей"""
    return render_template('dashboard.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


@app.route('/admin')
@login_required
@role_required_from_config()
def admin_panel():
    """Админ панель - только для админа"""
    return render_template('admin_panel.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


@app.route('/manager')
@login_required
@role_required_from_config()
def manager_panel():
    """Панель менеджера - для админа и менеджера"""
    return render_template('manager_panel.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


@app.route('/users')
@login_required
@role_required_from_config()
def users_list():
    """Управление пользователями - только для админа"""
    return render_template('users_list.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


# @app.route('/pcm-temp/<filename>')
# @login_required
# def serve_pcm_temp(filename):
#     import tempfile
#     return send_from_directory(tempfile.gettempdir(), filename)

@app.route('/settings')
@login_required
@role_required_from_config()
def settings():
    """Настройки - для админа и менеджера"""
    return render_template('settings.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


@app.route('/reports')
@login_required
@role_required_from_config()
def reports():
    """Отчеты - для всех кроме гостя"""
    return render_template('reports.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


# ============================================
# АУТЕНТИФИКАЦИЯ
# ============================================

@app.route('/login', methods=['GET', 'POST'])
@rate_limit(max_requests=10, window=60)
def login():
    """Страница входа"""
    if session.get('logged_in'):
        return redirect(url_for('index'))

    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')
        ip = request.remote_addr

        if not username or not password:
            flash('Введите имя пользователя и пароль.', 'danger')
            return render_template('login.html')

        # Проверяем попытки входа
        if not auth_manager.check_login_attempts(ip):
            flash('Слишком много попыток входа. Подождите 5 минут.', 'danger')
            logger.warning(f'Login blocked for IP: {ip}')
            return render_template('login.html')

        # Проверяем пользователя
        if username in user_hashes:
            if bcrypt.checkpw(password.encode('utf-8'), user_hashes[username]):
                # Успешный вход
                session['logged_in'] = True
                session['username'] = username
                session['user_role'] = user_roles.get(username, 'user')
                session['user_name'] = user_names.get(username, username)
                session['login_time'] = datetime.now().isoformat()

                # Сбрасываем попытки
                if ip in auth_manager.failed_attempts:
                    del auth_manager.failed_attempts[ip]
                if ip in auth_manager.locked_ips:
                    del auth_manager.locked_ips[ip]

                logger.info(f'Successful login from IP: {ip}, user: {username}, role: {session["user_role"]}')

                # Перенаправляем на страницу, с которой пришли, или на главную
                next_page = request.args.get('next')
                if next_page:
                    return redirect(next_page)
                return redirect(url_for('index'))

        # Неудачная попытка
        auth_manager.register_failed_attempt(ip)
        flash('Неверное имя пользователя или пароль.', 'danger')
        logger.warning(f'Failed login attempt from IP: {ip}, username: {username}')

        if ip in auth_manager.locked_ips:
            flash('IP-адрес заблокирован на 5 минут из-за множества неудачных попыток.', 'danger')

        return render_template('login.html')

    return render_template('login.html')


@app.route('/logout')
def logout():
    """Выход из системы"""
    session.clear()
    flash('Вы успешно вышли из системы.', 'info')
    return redirect(url_for('login'))


# ============================================
# API
# ============================================

@app.route('/api/coordinates')
@login_required
def get_coordinates():
    """API для получения координат"""
    return jsonify(coordinates)


@app.route('/api/users')
@login_required
@role_required(['admin'])
def get_users():
    """API для получения списка пользователей (только для админа)"""
    users = []
    for username, password in Config.USERS.items():
        users.append({
            'username': username,
            'role': user_roles.get(username, 'user'),
            'name': user_names.get(username, username)
        })
    return jsonify(users)


@app.route('/api/user/role')
@login_required
def get_user_role():
    """API для получения роли текущего пользователя"""
    return jsonify({
        'username': session.get('username'),
        'role': session.get('user_role'),
        'name': session.get('user_name')
    })


# ============================================
# ОБРАБОТКА ОШИБОК
# ============================================

@app.errorhandler(401)
def unauthorized(error):
    """Обработка неавторизованного доступа"""
    flash('Пожалуйста, авторизуйтесь для доступа к этой странице.', 'warning')
    return redirect(url_for('login'))


@app.errorhandler(403)
def forbidden(error):
    """Обработка запрещенного доступа"""
    return render_template('403.html'), 403


@app.errorhandler(404)
def not_found(error):
    """Обработка 404 ошибки"""
    return render_template('404.html'), 404


@app.errorhandler(429)
def too_many_requests(error):
    """Обработка слишком частых запросов"""
    return jsonify({'error': 'Слишком много запросов. Подождите немного.'}), 429


# ============================================
# ЗАПУСК
# ============================================

@app.route('/static/data/<path:filename>')
def serve_geojson(filename):
    """Сервинг GeoJSON файлов"""
    return send_from_directory('static/data', filename)

# app.py - добавляем новый маршрут



@app.route('/analytics')
@login_required
@role_required_from_config()
def analytics():
    """Страница аналитики - переворот строки"""
    return render_template('analytics.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'),
                           reversed_text=None,  # Передаем None, если еще ничего не вводили
                           original_text=None)


# ============================================
# API ДЛЯ АНАЛИТИКИ
# ============================================

@app.route('/api/reverse', methods=['POST'])
@login_required
def reverse_string():
    """
    API для переворота строки
    Принимает: {"text": "строка"}
    Возвращает: {"original": "строка", "reversed": "акортс"}
    """
    try:
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({'error': 'Поле "text" обязательно'}), 400

        original_text = data['text'].strip()
        if not original_text:
            return jsonify({'error': 'Текст не может быть пустым'}), 400

        # Переворачиваем строку
        reversed_text = original_text[::-1]

        # Логируем действие
        logger.info(f'User {session.get("username")} reversed string: "{original_text}" -> "{reversed_text}"')

        return jsonify({
            'success': True,
            'original': original_text,
            'reversed': reversed_text,
            'length': len(original_text)
        })

    except Exception as e:
        logger.error(f'Error in reverse_string: {str(e)}')
        return jsonify({'error': 'Произошла ошибка на сервере'}), 500


# ============================================
# СТРАНИЦА TEST - Конвертер температуры
# ============================================

@app.route('/test')
@login_required
@role_required_from_config()
def test_page():
    """Страница test - конвертер температуры"""
    return render_template('test.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'))


# ============================================
# API ДЛЯ СТРАНИЦЫ TEST - GET запросы
# ============================================

@app.route('/api/convert', methods=['GET'])
@login_required
def convert_temperature():
    """
    Конвертер температуры через GET
    Использование: /api/convert?value=25&from=C&to=F
    Параметры:
        value - значение температуры
        from - исходная единица (C или F)
        to - целевая единица (C или F)
    """
    try:
        # Получаем параметры из URL
        value_str = request.args.get('value')
        from_unit = request.args.get('from', 'C').upper()
        to_unit = request.args.get('to', 'F').upper()

        # Проверяем наличие значения
        if not value_str:
            return jsonify({'error': 'Параметр "value" обязателен'}), 400

        # Конвертируем значение в число
        try:
            value = float(value_str)
        except ValueError:
            return jsonify({'error': 'Значение должно быть числом'}), 400

        # Проверяем единицы измерения
        if from_unit not in ['C', 'F']:
            return jsonify({'error': 'Параметр "from" должен быть C или F'}), 400

        if to_unit not in ['C', 'F']:
            return jsonify({'error': 'Параметр "to" должен быть C или F'}), 400

        # Если единицы одинаковые - возвращаем то же значение
        if from_unit == to_unit:
            result = value
            formula = f"{value}°{from_unit} = {value}°{to_unit} (без изменений)"
        else:
            # Конвертация
            if from_unit == 'C' and to_unit == 'F':
                result = (value * 9 / 5) + 32
                formula = f"({value} × 9/5) + 32 = {result:.2f}°F"
            elif from_unit == 'F' and to_unit == 'C':
                result = (value - 32) * 5 / 9
                formula = f"({value} - 32) × 5/9 = {result:.2f}°C"
            else:
                return jsonify({'error': 'Неверные единицы измерения'}), 400

        # Логируем действие
        logger.info(
            f'User {session.get("username")} converted temperature: {value}°{from_unit} -> {result:.2f}°{to_unit}')

        return jsonify({
            'success': True,
            'original': {
                'value': value,
                'unit': from_unit
            },
            'result': {
                'value': round(result, 2),
                'unit': to_unit
            },
            'formula': formula
        })

    except Exception as e:
        logger.error(f'Error in convert_temperature: {str(e)}')
        return jsonify({'error': 'Произошла ошибка на сервере'}), 500


@app.route('/api/weather', methods=['GET'])
@login_required
def get_weather_info():
    """
    Пример GET API с несколькими параметрами
    Использование: /api/weather?city=Moscow&units=metric
    """
    city = request.args.get('city', 'Moscow')
    units = request.args.get('units', 'metric')

    # Простая заглушка с данными
    weather_data = {
        'Moscow': {'temp': 5, 'condition': 'Облачно', 'wind': '3 м/с'},
        'London': {'temp': 8, 'condition': 'Дождь', 'wind': '5 м/с'},
        'Paris': {'temp': 10, 'condition': 'Солнечно', 'wind': '2 м/с'},
        'Berlin': {'temp': 6, 'condition': 'Туман', 'wind': '4 м/с'},
        'Madrid': {'temp': 15, 'condition': 'Ясно', 'wind': '1 м/с'},
    }

    # Ищем город (регистронезависимо)
    for key in weather_data:
        if key.lower() == city.lower():
            data = weather_data[key]
            temp = data['temp']
            if units == 'fahrenheit':
                temp = (temp * 9 / 5) + 32

            return jsonify({
                'success': True,
                'city': key,
                'temperature': round(temp, 1),
                'units': units,
                'condition': data['condition'],
                'wind': data['wind']
            })

    return jsonify({
        'success': False,
        'error': f'Город "{city}" не найден',
        'available_cities': list(weather_data.keys())
    }), 404


# ============================================
# СТРАНИЦА ДОКУМЕНТОВ
# ============================================

@app.route('/documents')
@login_required
@role_required_from_config()
def documents():
    """Страница с таблицей документов - с навигацией по папкам"""
    import os
    import json
    import striprtf
    # Базовая папка с документами
    base_folder = os.path.join(app.root_path, 'documents')

    # Получаем текущий путь из параметра
    current_subpath = request.args.get('path', '')

    # Формируем полный путь
    if current_subpath:
        current_path = os.path.normpath(os.path.join(base_folder, current_subpath))
        if not current_path.startswith(os.path.normpath(base_folder)):
            current_path = base_folder
            current_subpath = ''
    else:
        current_path = base_folder

    # Проверяем существование папки
    if not os.path.exists(current_path):
        os.makedirs(current_path, exist_ok=True)
        logger.info(f'Created folder: {current_path}')

    # Собираем папки и файлы
    folders = []
    files = []

    supported_types = {
        'pdf': 'pdf',
        'doc': 'docx',
        'docx': 'docx',
        'rtf':'docx',
        'wav': 'wav',
        'mp3': 'audio',
        'ogg': 'audio',
        'flac': 'audio',
        'm4a': 'audio',
        'txt': 'txt',
        'log': 'txt',
        'csv': 'txt',
        'md': 'txt',
        'py': 'txt',
        'js': 'txt',
        'html': 'txt',
        'css': 'txt',
        'json': 'txt',
        'xml': 'txt',
        'xls': 'excel',
        'xlsx': 'excel',
        'ppt': 'pptx',
        'pptx': 'pptx',
        'jpg': 'image',
        'jpeg': 'image',
        'png': 'image',
        'gif': 'image',
        'bmp': 'image',
        'webp': 'image',
        'svg': 'image',
        'mp4': 'video',
        'webm': 'video',
        'avi': 'video',
        'mov': 'video',
        'mkv': 'video'
    }

    try:
        items = os.listdir(current_path)
        items.sort(key=lambda x: (not os.path.isdir(os.path.join(current_path, x)), x.lower()))

        for item in items:
            item_path = os.path.join(current_path, item)

            if os.path.isdir(item_path):
                # Это папка
                subpath = os.path.join(current_subpath, item).replace('\\', '/') if current_subpath else item
                # Считаем файлы внутри (рекурсивно)
                file_count = count_files_recursive(item_path)
                folders.append({
                    'name': item,
                    'path': subpath,
                    'file_count': file_count
                })
            else:
                # Это файл
                file_extension = item.lower().split('.')[-1] if '.' in item else 'unknown'
                file_type = supported_types.get(file_extension, 'other')

                # Размер файла
                try:
                    file_size = os.path.getsize(item_path)
                    if file_size < 1024:
                        size_str = f"{file_size} B"
                    elif file_size < 1024 * 1024:
                        size_str = f"{file_size / 1024:.1f} KB"
                    elif file_size < 1024 * 1024 * 1024:
                        size_str = f"{file_size / (1024 * 1024):.1f} MB"
                    else:
                        size_str = f"{file_size / (1024 * 1024 * 1024):.1f} GB"
                except:
                    size_str = "N/A"

                # Дата модификации
                try:
                    mtime = os.path.getmtime(item_path)
                    date_modified = datetime.fromtimestamp(mtime).strftime('%Y-%m-%d %H:%M')
                except:
                    date_modified = "Неизвестно"

                # Путь относительно базовой папки
                file_rel_path = os.path.join(current_subpath, item).replace('\\', '/') if current_subpath else item

                # Имя без расширения
                name_without_ext = os.path.splitext(item)[0]
                display_name = name_without_ext

                files.append({
                    'path': file_rel_path,
                    'title': display_name,
                    'date': date_modified,
                    'type': file_type,
                    'extension': file_extension,
                    'size': size_str
                })

        logger.info(f'Found {len(folders)} folders and {len(files)} files in {current_path}')

    except Exception as e:
        logger.error(f'Error scanning folder: {str(e)}')
        flash(f'Ошибка при сканировании папки: {str(e)}', 'danger')

    # Хлебные крошки
    breadcrumbs = [{'name': '📁 Все документы', 'path': ''}]
    if current_subpath:
        parts = current_subpath.replace('\\', '/').split('/')
        accumulated = ''
        for part in parts:
            accumulated = os.path.join(accumulated, part).replace('\\', '/') if accumulated else part
            breadcrumbs.append({'name': part, 'path': accumulated})

    # JSON для JavaScript
    files_json = json.dumps(files, ensure_ascii=False)

    return render_template('documents.html',
                           username=session.get('username'),
                           user_role=session.get('user_role'),
                           user_name=session.get('user_name'),
                           folders=folders,
                           files=files,
                           files_json=files_json,
                           current_path=current_subpath,
                           breadcrumbs=breadcrumbs)


def count_files_recursive(folder_path):
    import os
    """Рекурсивный подсчет файлов в папке"""
    count = 0
    try:
        for item in os.listdir(folder_path):
            item_path = os.path.join(folder_path, item)
            if os.path.isfile(item_path):
                count += 1
            elif os.path.isdir(item_path):
                count += count_files_recursive(item_path)
    except:
        pass
    return count


# ============================================
# API ДЛЯ ПРОСМОТРА ДОКУМЕНТОВ
# ============================================

@app.route('/api/view-document/<filename>')
@login_required
def view_document(filename):
    """
    API для просмотра документа в новой вкладке
    Возвращает HTML страницу с предпросмотром или плеером
    """
    import os
    from flask import make_response, send_file
    import html as html_module
    import striprtf

    # Путь к папке с документами
    documents_folder = os.path.join(app.root_path, 'documents')
    file_path = os.path.join(documents_folder, filename)

    # Проверяем существование файла
    if not os.path.exists(file_path):
        error_html = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Ошибка - Файл не найден</title>
            <style>
                body {{
                    font-family: Arial, sans-serif;
                    display: flex;
                    justify-content: center;
                    align-items: center;
                    height: 100vh;
                    margin: 0;
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                }}
                .error-container {{
                    text-align: center;
                    background: white;
                    padding: 40px;
                    border-radius: 10px;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.3);
                    max-width: 500px;
                }}
                .error-icon {{
                    font-size: 64px;
                    margin-bottom: 20px;
                }}
                h1 {{
                    color: #e74c3c;
                    margin-bottom: 15px;
                }}
                p {{
                    color: #666;
                    margin-bottom: 20px;
                    line-height: 1.6;
                }}
                .btn {{
                    display: inline-block;
                    padding: 10px 20px;
                    background: #667eea;
                    color: white;
                    text-decoration: none;
                    border-radius: 5px;
                    transition: background 0.3s;
                    margin: 5px;
                }}
                .btn:hover {{
                    background: #764ba2;
                }}
            </style>
        </head>
        <body>
            <div class="error-container">
                <div class="error-icon">📄❌</div>
                <h1>Ошибка: Файл не найден</h1>
                <p>Файл <strong>"{filename}"</strong> не существует или был удален.<br>
                Проверьте правильность имени файла или обратитесь к администратору.</p>
                <a href="javascript:window.close()" class="btn">Закрыть вкладку</a>
            </div>
        </body>
        </html>
        '''
        return make_response(error_html, 404)

    # Определяем тип файла по расширению
    file_extension = filename.lower().split('.')[-1] if '.' in filename else ''

    # Проверяем, не битый ли файл
    try:
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            raise ValueError("Файл пустой")
    except Exception as e:
        error_html = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Ошибка - Битый файл</title>
            <style>
                body {{
                    font-family: Arial, sans-serif;
                    display: flex;
                    justify-content: center;
                    align-items: center;
                    height: 100vh;
                    margin: 0;
                    background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                }}
                .error-container {{
                    text-align: center;
                    background: white;
                    padding: 40px;
                    border-radius: 10px;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.3);
                    max-width: 500px;
                }}
                .error-icon {{
                    font-size: 64px;
                    margin-bottom: 20px;
                }}
                h1 {{
                    color: #e74c3c;
                    margin-bottom: 15px;
                }}
                p {{
                    color: #666;
                    margin-bottom: 20px;
                    line-height: 1.6;
                }}
                .btn {{
                    display: inline-block;
                    padding: 10px 20px;
                    background: #667eea;
                    color: white;
                    text-decoration: none;
                    border-radius: 5px;
                    transition: background 0.3s;
                    margin: 5px;
                }}
                .btn:hover {{
                    background: #764ba2;
                }}
                .error-details {{
                    background: #f8f9fa;
                    padding: 15px;
                    border-radius: 5px;
                    margin: 15px 0;
                    font-size: 14px;
                    color: #e74c3c;
                }}
            </style>
        </head>
        <body>
            <div class="error-container">
                <div class="error-icon">⚠️💥</div>
                <h1>Ошибка: Файл поврежден</h1>
                <p>Файл <strong>"{filename}"</strong> поврежден или имеет некорректный формат.<br>
                Невозможно отобразить содержимое файла.</p>
                <div class="error-details">
                    <strong>Причина:</strong> {str(e)}
                </div>
                <a href="javascript:window.close()" class="btn">Закрыть вкладку</a>
            </div>
        </body>
        </html>
        '''
        return make_response(error_html, 500)

    # ============================================
    # ДЛЯ ИЗОБРАЖЕНИЙ (jpg, jpeg, png, gif, bmp, webp)
    # ============================================
    if file_extension in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'svg', 'ico']:
        # Определяем MIME-тип
        mime_types = {
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'gif': 'image/gif',
            'bmp': 'image/bmp',
            'webp': 'image/webp',
            'svg': 'image/svg+xml',
            'ico': 'image/x-icon'
        }
        mime_type = mime_types.get(file_extension, 'image/jpeg')

        # Проверяем, не битое ли изображение
        try:
            from PIL import Image
            img = Image.open(file_path)
            img.verify()  # Проверяем целостность
            img = Image.open(file_path)  # Открываем заново после verify()
            width, height = img.size
            img_format = img.format
            img_mode = img.mode
            image_info = f"Размер: {width}x{height}px, Формат: {img_format}, Режим: {img_mode}"
        except:
            image_info = "Информация недоступна"

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Просмотр изображения - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{
                    background: #1a1a1a;
                    font-family: Arial, sans-serif;
                    display: flex;
                    flex-direction: column;
                    height: 100vh;
                }}
                .toolbar {{
                    background: #2c2c2c;
                    color: white;
                    padding: 10px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.5);
                    z-index: 10;
                }}
                .toolbar h3 {{
                    margin: 0;
                    font-size: 16px;
                    color: #fff;
                }}
                .toolbar .info {{
                    font-size: 12px;
                    color: #aaa;
                    margin-left: 20px;
                }}
                .btn {{
                    padding: 8px 15px;
                    background: #667eea;
                    color: white;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                    margin-left: 10px;
                }}
                .btn:hover {{ background: #764ba2; }}
                .image-container {{
                    flex: 1;
                    display: flex;
                    justify-content: center;
                    align-items: center;
                    overflow: auto;
                    padding: 20px;
                }}
                img {{
                    max-width: 95%;
                    max-height: 95vh;
                    object-fit: contain;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.5);
                    border-radius: 5px;
                    background: #2c2c2c;
                }}
                .error-message {{
                    display: none;
                    text-align: center;
                    color: white;
                    padding: 50px;
                }}
                .zoom-controls {{
                    display: flex;
                    gap: 10px;
                    align-items: center;
                }}
                .zoom-btn {{
                    background: #444;
                    color: white;
                    border: none;
                    width: 30px;
                    height: 30px;
                    border-radius: 50%;
                    cursor: pointer;
                    font-size: 18px;
                    display: flex;
                    align-items: center;
                    justify-content: center;
                }}
                .zoom-btn:hover {{ background: #667eea; }}
                .zoom-level {{
                    color: #aaa;
                    font-size: 12px;
                    min-width: 50px;
                    text-align: center;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <div style="display: flex; align-items: center;">
                    <h3>🖼️ {filename}</h3>
                    <span class="info">{image_info}</span>
                </div>
                <div style="display: flex; align-items: center;">
                    <div class="zoom-controls">
                        <button class="zoom-btn" onclick="zoomOut()">−</button>
                        <span class="zoom-level" id="zoomLevel">100%</span>
                        <button class="zoom-btn" onclick="zoomIn()">+</button>
                        <button class="zoom-btn" onclick="zoomReset()">↺</button>
                    </div>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="image-container" id="imageContainer">
                <img src="/documents/{filename}" alt="{filename}" id="mainImage" 
                     onerror="showError()" onload="hideError()">
            </div>
            <div class="error-message" id="errorMessage">
                <h2>❌ Не удалось загрузить изображение</h2>
                <p>Возможно, файл поврежден или имеет неверный формат.</p>
                <button onclick="window.close()" class="btn">Закрыть</button>
            </div>
            <script>
                let currentZoom = 1;
                const img = document.getElementById('mainImage');

                function updateZoom() {{
                    img.style.transform = `scale(${{currentZoom}})`;
                    document.getElementById('zoomLevel').textContent = Math.round(currentZoom * 100) + '%';
                }}

                function zoomIn() {{
                    if (currentZoom < 3) {{
                        currentZoom += 0.25;
                        updateZoom();
                    }}
                }}

                function zoomOut() {{
                    if (currentZoom > 0.25) {{
                        currentZoom -= 0.25;
                        updateZoom();
                    }}
                }}

                function zoomReset() {{
                    currentZoom = 1;
                    updateZoom();
                }}

                function showError() {{
                    document.getElementById('imageContainer').style.display = 'none';
                    document.getElementById('errorMessage').style.display = 'block';
                }}

                function hideError() {{
                    document.getElementById('imageContainer').style.display = 'flex';
                    document.getElementById('errorMessage').style.display = 'none';
                }}

                // Зум колесиком мыши
                document.addEventListener('wheel', function(e) {{
                    if (e.ctrlKey || e.metaKey) {{
                        e.preventDefault();
                        if (e.deltaY < 0) zoomIn();
                        else zoomOut();
                    }}
                }});
            </script>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ ТЕКСТОВЫХ ФАЙЛОВ (txt, log, csv, md, py, js, html, css, xml, json)
    # ============================================
    elif file_extension in ['txt', 'log', 'csv', 'md', 'py', 'js', 'html', 'css', 'xml', 'json', 'yaml', 'yml', 'ini',
                            'cfg', 'conf', 'sh', 'bat', 'sql', 'java', 'cpp', 'c', 'h', 'php', 'rb', 'go', 'rs',
                            'swift', 'kt', 'ts', 'tsx', 'jsx', 'vue', 'env', 'gitignore', 'dockerfile']:
        try:
            # Пробуем прочитать файл
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='cp1251') as f:
                    content = f.read()
            except:
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        content = f.read()
                except Exception as e:
                    error_html = f'''
                    <!DOCTYPE html>
                    <html lang="ru">
                    <head>
                        <meta charset="UTF-8">
                        <title>Ошибка чтения файла</title>
                        <style>
                            body {{
                                font-family: Arial, sans-serif;
                                display: flex;
                                justify-content: center;
                                align-items: center;
                                height: 100vh;
                                margin: 0;
                                background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                            }}
                            .error-container {{
                                text-align: center;
                                background: white;
                                padding: 40px;
                                border-radius: 10px;
                                box-shadow: 0 10px 30px rgba(0,0,0,0.3);
                                max-width: 500px;
                            }}
                            h1 {{ color: #e74c3c; }}
                            .btn {{
                                padding: 10px 20px;
                                background: #667eea;
                                color: white;
                                border: none;
                                border-radius: 5px;
                                cursor: pointer;
                                text-decoration: none;
                            }}
                        </style>
                    </head>
                    <body>
                        <div class="error-container">
                            <h1>⚠️ Ошибка чтения файла</h1>
                            <p>Не удалось прочитать файл <strong>"{filename}"</strong>.</p>
                            <p>Файл может быть бинарным или поврежден.</p>
                            <p style="color: #e74c3c;">{str(e)}</p>
                            <a href="/documents/{filename}" download class="btn">💾 Скачать файл</a>
                            <button onclick="window.close()" class="btn">Закрыть</button>
                        </div>
                    </body>
                    </html>
                    '''
                    return make_response(error_html, 500)
        except Exception as e:
            error_html = f'''
            <!DOCTYPE html>
            <html lang="ru">
            <head>
                <meta charset="UTF-8">
                <title>Ошибка чтения файла</title>
                <style>
                    body {{
                        font-family: Arial, sans-serif;
                        display: flex;
                        justify-content: center;
                        align-items: center;
                        height: 100vh;
                        margin: 0;
                        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
                    }}
                    .error-container {{
                        text-align: center;
                        background: white;
                        padding: 40px;
                        border-radius: 10px;
                        box-shadow: 0 10px 30px rgba(0,0,0,0.3);
                        max-width: 500px;
                    }}
                    h1 {{ color: #e74c3c; }}
                    .btn {{
                        padding: 10px 20px;
                        background: #667eea;
                        color: white;
                        border: none;
                        border-radius: 5px;
                        cursor: pointer;
                        text-decoration: none;
                    }}
                </style>
            </head>
            <body>
                <div class="error-container">
                    <h1>⚠️ Ошибка чтения файла</h1>
                    <p>Не удалось прочитать файл <strong>"{filename}"</strong>.</p>
                    <p style="color: #e74c3c;">{str(e)}</p>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать файл</a>
                    <button onclick="window.close()" class="btn">Закрыть</button>
                </div>
            </body>
            </html>
            '''
            return make_response(error_html, 500)

        # Экранируем HTML
        escaped_content = html_module.escape(content)

        # Подсвечиваем синтаксис (простая версия)
        lines = escaped_content.split('\n')
        line_count = len(lines)

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Просмотр текста - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{
                    background: #1e1e1e;
                    font-family: 'Consolas', 'Monaco', 'Courier New', monospace;
                    display: flex;
                    flex-direction: column;
                    height: 100vh;
                }}
                .toolbar {{
                    background: #2d2d2d;
                    color: #d4d4d4;
                    padding: 10px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.5);
                    flex-wrap: wrap;
                    gap: 10px;
                }}
                .toolbar h3 {{
                    margin: 0;
                    font-size: 16px;
                }}
                .file-info {{
                    font-size: 12px;
                    color: #888;
                }}
                .btn {{
                    padding: 8px 15px;
                    background: #0e639c;
                    color: white;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                }}
                .btn:hover {{ background: #1177bb; }}
                .content-area {{
                    flex: 1;
                    overflow: auto;
                    padding: 15px;
                    background: #1e1e1e;
                }}
                pre {{
                    margin: 0;
                    color: #d4d4d4;
                    font-size: 14px;
                    line-height: 1.6;
                    white-space: pre-wrap;
                    word-wrap: break-word;
                }}
                .line-numbers {{
                    color: #858585;
                    user-select: none;
                    padding-right: 20px;
                    border-right: 1px solid #404040;
                    margin-right: 20px;
                    display: inline-block;
                    text-align: right;
                    min-width: 40px;
                }}
                .code-line {{
                    display: flex;
                }}
                .code-line:hover {{
                    background: #2a2d2e;
                }}
                .error-message {{
                    display: none;
                    text-align: center;
                    color: #d4d4d4;
                    padding: 50px;
                }}
                ::-webkit-scrollbar {{
                    width: 10px;
                    height: 10px;
                }}
                ::-webkit-scrollbar-track {{
                    background: #1e1e1e;
                }}
                ::-webkit-scrollbar-thumb {{
                    background: #424242;
                    border-radius: 5px;
                }}
                ::-webkit-scrollbar-thumb:hover {{
                    background: #4e4e4e;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <div>
                    <h3>📄 {filename}</h3>
                    <span class="file-info">Строк: {line_count} | Размер: {file_size} байт</span>
                </div>
                <div>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="content-area" id="contentArea">
                <pre id="textContent">{escaped_content}</pre>
            </div>
            <div class="error-message" id="errorMessage">
                <h2>❌ Ошибка отображения</h2>
                <p>Не удалось отобразить содержимое файла.</p>
                <button onclick="window.close()" class="btn">Закрыть</button>
            </div>
        </body>
        </html>
        '''
        return make_response(html_content)

    elif file_extension == 'rtf':
        text_preview = ""
        try:
            from striprtf.striprtf import rtf_to_text

            # Читаем файл
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()

            # Конвертируем в текст
            text_preview = rtf_to_text(rtf_content)

            if not text_preview or not text_preview.strip():
                text_preview = "(документ пустой)"

        except ImportError:
            text_preview = "Установите: pip install striprtf"
        except Exception as e:
            text_preview = f"Ошибка: {str(e)}"

        escaped_text = html_module.escape(text_preview)

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <title>Просмотр RTF - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{ 
                    background: #f5f5f5; 
                    font-family: 'Segoe UI', Arial, sans-serif; 
                    display: flex; 
                    flex-direction: column; 
                    height: 100vh; 
                }}
                .toolbar {{
                    background: #6c5ce7;
                    color: white;
                    padding: 15px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                }}
                .toolbar h3 {{ margin: 0; font-size: 18px; }}
                .btn {{
                    padding: 8px 15px;
                    background: white;
                    color: #6c5ce7;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-weight: bold;
                    font-size: 14px;
                }}
                .btn:hover {{ background: #e0e0e0; }}
                .content {{
                    flex: 1;
                    overflow-y: auto;
                    padding: 30px 40px;
                    background: white;
                    margin: 20px;
                    border-radius: 10px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                }}
                pre {{
                    white-space: pre-wrap;
                    word-wrap: break-word;
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 16px;
                    line-height: 1.8;
                    color: #333;
                }}
                .info {{
                    color: #666;
                    font-size: 14px;
                    margin-bottom: 20px;
                    padding-bottom: 15px;
                    border-bottom: 2px solid #6c5ce7;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <h3>📝 {filename}</h3>
                <div style="display: flex; gap: 10px;">
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="content">
                <div class="info">
                    📄 Файл: <strong>{filename}</strong> | 
                    📏 Размер: <strong>{file_size} байт</strong>
                </div>
                <pre>{escaped_text}</pre>
            </div>
        </body>
        </html>
        '''
        return make_response(html_content)

    elif file_extension == 'pcm':
        import struct

        # === ПАРАМЕТРЫ PCM (пропиши свои) ===
        PCM_RATE = 6000  # частота дискретизации
        PCM_CHANNELS = 1  # моно
        PCM_BITS = 16  # 16-bit

        try:
            # Читаем PCM данные
            with open(file_path, 'rb') as f:
                pcm_data = f.read()

            # Создаем WAV в памяти
            wav_data = io.BytesIO()

            # Размер данных
            data_size = len(pcm_data)
            sample_width = PCM_BITS // 8  # байт на сэмпл
            byte_rate = PCM_RATE * PCM_CHANNELS * sample_width

            # WAV заголовок (44 байта)
            wav_data.write(b'RIFF')
            wav_data.write(struct.pack('<I', 36 + data_size))  # размер файла - 8
            wav_data.write(b'WAVE')

            # fmt chunk
            wav_data.write(b'fmt ')
            wav_data.write(struct.pack('<I', 16))  # размер fmt
            wav_data.write(struct.pack('<H', 1))  # PCM = 1
            wav_data.write(struct.pack('<H', PCM_CHANNELS))  # каналы
            wav_data.write(struct.pack('<I', PCM_RATE))  # частота
            wav_data.write(struct.pack('<I', byte_rate))  # байт/сек
            wav_data.write(struct.pack('<H', PCM_CHANNELS * sample_width))  # блок
            wav_data.write(struct.pack('<H', PCM_BITS))  # бит на сэмпл

            # data chunk
            wav_data.write(b'data')
            wav_data.write(struct.pack('<I', data_size))  # размер данных
            wav_data.write(pcm_data)  # сами данные

            wav_data.seek(0)

            # Сохраняем временный WAV
            wav_temp = tempfile.NamedTemporaryFile(suffix='.wav', delete=False)
            wav_temp.write(wav_data.read())
            wav_path = wav_temp.name
            wav_temp.close()

            # Показываем плеер
            html_content = f'''
            <!DOCTYPE html>
            <html lang="ru">
            <head>
                <meta charset="UTF-8">
                <title>Аудио - {filename}</title>
                <style>
                    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                    body {{
                        font-family: Arial, sans-serif;
                        background: linear-gradient(135deg, #1e3c72 0%, #2a5298 100%);
                        display: flex;
                        justify-content: center;
                        align-items: center;
                        height: 100vh;
                    }}
                    .player-container {{
                        background: white;
                        padding: 40px;
                        border-radius: 20px;
                        box-shadow: 0 20px 60px rgba(0,0,0,0.3);
                        text-align: center;
                        max-width: 600px;
                        width: 90%;
                    }}
                    .audio-icon {{ font-size: 80px; margin-bottom: 20px; }}
                    h2 {{ color: #333; margin-bottom: 10px; }}
                    .info {{ color: #666; font-size: 14px; margin-bottom: 20px; }}
                    audio {{ width: 100%; margin-bottom: 20px; }}
                    .btn {{
                        padding: 10px 20px;
                        background: #667eea;
                        color: white;
                        border: none;
                        border-radius: 8px;
                        cursor: pointer;
                        text-decoration: none;
                        font-size: 14px;
                        margin: 5px;
                        display: inline-block;
                    }}
                    .btn:hover {{ background: #764ba2; }}
                </style>
            </head>
            <body>
                <div class="player-container">
                    <div class="audio-icon">🎵</div>
                    <h2>{filename}</h2>
                    <div class="info">
                        {PCM_RATE} Гц | {'Моно' if PCM_CHANNELS == 1 else 'Стерео'} | {PCM_BITS}-bit
                    </div>
                    <audio controls autoplay>
                        <source src="/pcm-temp/{os.path.basename(wav_path)}" type="audio/wav">
                    </audio>
                    <div>
                        <a href="/documents/{filename}" download class="btn">💾 Скачать PCM</a>
                        <button onclick="window.close()" class="btn">✖ Закрыть</button>
                    </div>
                </div>
            </body>
            </html>
            '''
            return make_response(html_content)

        except Exception as e:
            html_content = f'''
            <!DOCTYPE html>
            <html lang="ru">
            <head>
                <meta charset="UTF-8">
                <title>PCM файл - {filename}</title>
                <style>
                    body {{ 
                        font-family: Arial; 
                        display: flex; 
                        justify-content: center; 
                        align-items: center; 
                        height: 100vh; 
                        background: #f5f5f5; 
                    }}
                    .box {{ 
                        text-align: center; 
                        background: white; 
                        padding: 40px; 
                        border-radius: 15px; 
                    }}
                    .btn {{ 
                        padding: 10px 20px; 
                        background: #667eea; 
                        color: white; 
                        text-decoration: none; 
                        border-radius: 8px; 
                    }}
                </style>
            </head>
            <body>
                <div class="box">
                    <h2>🎵 {filename}</h2>
                    <p>Ошибка: {str(e)}</p>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                </div>
            </body>
            </html>
            '''
            return make_response(html_content)
    # ============================================
    # ДЛЯ PDF ФАЙЛОВ
    # ============================================
    elif file_extension == 'pdf':
        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Просмотр PDF - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{ background: #525659; font-family: Arial, sans-serif; }}
                .toolbar {{
                    background: #323639;
                    color: white;
                    padding: 10px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.3);
                }}
                .toolbar h3 {{ margin: 0; font-size: 16px; }}
                .btn {{
                    padding: 8px 15px;
                    background: #667eea;
                    color: white;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                }}
                .btn:hover {{ background: #764ba2; }}
                embed {{
                    width: 100%;
                    height: calc(100vh - 50px);
                    border: none;
                }}
                .error-message {{
                    display: none;
                    text-align: center;
                    padding: 50px;
                    color: white;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <h3>📕 {filename}</h3>
                <div>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <embed src="/documents/{filename}" type="application/pdf" id="pdfViewer">
            <div class="error-message" id="errorMessage">
                <h2>Не удалось загрузить PDF</h2>
                <p>Возможно, файл поврежден или имеет неверный формат.</p>
                <button onclick="window.close()" class="btn">Закрыть</button>
            </div>
            <script>
                document.getElementById('pdfViewer').addEventListener('error', function() {{
                    document.getElementById('pdfViewer').style.display = 'none';
                    document.getElementById('errorMessage').style.display = 'block';
                }});
            </script>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ АУДИО ФАЙЛОВ (wav, mp3, ogg, flac, m4a)
    # ============================================
    elif file_extension in ['wav', 'mp3', 'ogg', 'flac', 'm4a', 'aac', 'wma']:
        # Определяем MIME-тип
        mime_types = {
            'wav': 'audio/wav',
            'mp3': 'audio/mpeg',
            'ogg': 'audio/ogg',
            'flac': 'audio/flac',
            'm4a': 'audio/mp4',
            'aac': 'audio/aac',
            'wma': 'audio/x-ms-wma'
        }
        mime_type = mime_types.get(file_extension, 'audio/wav')

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Аудио плеер - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{
                    font-family: Arial, sans-serif;
                    background: linear-gradient(135deg, #1e3c72 0%, #2a5298 100%);
                    display: flex;
                    justify-content: center;
                    align-items: center;
                    height: 100vh;
                }}
                .player-container {{
                    background: white;
                    padding: 40px;
                    border-radius: 20px;
                    box-shadow: 0 20px 60px rgba(0,0,0,0.3);
                    text-align: center;
                    max-width: 600px;
                    width: 90%;
                }}
                .audio-icon {{
                    font-size: 80px;
                    margin-bottom: 20px;
                    animation: pulse 2s infinite;
                }}
                @keyframes pulse {{
                    0%, 100% {{ transform: scale(1); }}
                    50% {{ transform: scale(1.1); }}
                }}
                h2 {{
                    color: #333;
                    margin-bottom: 10px;
                    word-wrap: break-word;
                }}
                .file-info {{
                    color: #666;
                    margin-bottom: 30px;
                    font-size: 14px;
                }}
                audio {{
                    width: 100%;
                    margin-bottom: 20px;
                }}
                .controls {{
                    display: flex;
                    gap: 10px;
                    justify-content: center;
                    flex-wrap: wrap;
                }}
                .btn {{
                    padding: 10px 20px;
                    background: #667eea;
                    color: white;
                    border: none;
                    border-radius: 8px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                    transition: background 0.3s;
                }}
                .btn:hover {{ background: #764ba2; }}
                .error-message {{
                    display: none;
                    color: #e74c3c;
                    margin-top: 15px;
                }}
            </style>
        </head>
        <body>
            <div class="player-container">
                <div class="audio-icon">🎵</div>
                <h2>Аудио файл</h2>
                <p class="file-info">{filename}</p>
                <audio controls autoplay id="audioPlayer">
                    <source src="/documents/{filename}" type="{mime_type}">
                    Ваш браузер не поддерживает аудио элемент.
                </audio>
                <div class="error-message" id="errorMessage">
                    <p>⚠️ Ошибка воспроизведения. Файл может быть поврежден.</p>
                </div>
                <div class="controls">
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <script>
                var audio = document.getElementById('audioPlayer');
                audio.addEventListener('error', function() {{
                    document.getElementById('errorMessage').style.display = 'block';
                }});
                audio.addEventListener('loadeddata', function() {{
                    document.getElementById('errorMessage').style.display = 'none';
                }});
            </script>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ WORD ДОКУМЕНТОВ (docx, doc) - ЛОКАЛЬНАЯ СЕТЬ
    # ============================================
    elif file_extension in ['docx', 'doc']:
        # Извлекаем текст
        text_preview = ""
        recovery_needed = False
        error_msg = ""
        recovery_method = ""

        try:
            if file_extension == 'docx':
                try:
                    from docx import Document
                    doc = Document(file_path)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    text_preview = '\n'.join(paragraphs)
                    if not text_preview:
                        text_preview = "(документ не содержит текста)"
                except ImportError:
                    text_preview = "⚠️ Библиотека python-docx не установлена. Текст не извлечен.\nУстановите: pip install python-docx"
                    recovery_needed = True
                    recovery_method = "library_missing"
                except Exception as e:
                    error_msg = str(e)
                    recovery_needed = True
                    recovery_method = "read_error"
                    # Пробуем восстановить через zipfile
                    try:
                        import zipfile
                        import re

                        with zipfile.ZipFile(file_path, 'r') as z:
                            if 'word/document.xml' in z.namelist():
                                xml_content = z.read('word/document.xml')
                                texts = re.findall(r'<w:t[^>]*>(.*?)</w:t>',
                                                   xml_content.decode('utf-8', errors='ignore'))
                                recovered_text = [t for t in texts if t.strip()]

                                if recovered_text:
                                    text_preview = '\n'.join(recovered_text)
                                    recovery_method = "xml_extracted"
                                else:
                                    text_preview = "⚠️ Не удалось извлечь текст из файла."
                                    recovery_method = "xml_empty"
                            else:
                                text_preview = "⚠️ Файл поврежден: отсутствует document.xml."
                                recovery_method = "no_xml"
                    except zipfile.BadZipFile:
                        text_preview = "⚠️ Файл серьезно поврежден: не является ZIP-архивом."
                        recovery_method = "bad_zip"
                    except Exception as recovery_error:
                        text_preview = f"⚠️ Не удалось восстановить файл.\nОшибка: {str(recovery_error)}"
                        recovery_method = "recovery_failed"
            else:
                # .doc файлы
                text_preview = "⚠️ Файлы .doc (старый формат) не поддерживаются для просмотра."
                recovery_needed = True
                recovery_method = "old_format"
        except Exception as e:
            text_preview = f"❌ Критическая ошибка: {str(e)}"
            recovery_needed = True
            recovery_method = "critical_error"

        escaped_text = html_module.escape(text_preview)

        # Определяем статус и сообщение для пользователя
        if not recovery_needed:
            status_badge = '<span class="status-badge status-ok">✅ OK</span>'
            status_title = "Документ открыт успешно"
        elif recovery_method in ["xml_extracted"]:
            status_badge = '<span class="status-badge status-recovered">⚠️ Частично восстановлен</span>'
            status_title = "Текст извлечен, но документ поврежден"
        else:
            status_badge = '<span class="status-badge status-error">❌ Поврежден</span>'
            status_title = "Документ поврежден"

        # Текст в зависимости от ситуации
        if recovery_method in ["read_error", "xml_extracted", "xml_empty", "no_xml", "bad_zip", "recovery_failed"]:
            recovery_text = """
            <h4>🔧 Документ поврежден — требуется восстановление в Word</h4>
            <p>Файл имеет поврежденную структуру. <strong>Microsoft Word</strong> при открытии автоматически 
            пытается восстановить поврежденные документы. Для восстановления:</p>
            <ol style="margin: 10px 0 15px 20px; color: #856404; line-height: 1.8;">
                <li>Скачайте файл на ваш компьютер</li>
                <li>Откройте его в Microsoft Word</li>
                <li>Word предложит восстановить документ — нажмите "Восстановить"</li>
                <li>Сохраните восстановленный файл</li>
            </ol>
            """
        elif recovery_method == "old_format":
            recovery_text = """
            <h4>📜 Старый формат .doc</h4>
            <p>Файлы в формате <strong>.doc</strong> (Microsoft Word 97-2003) не поддерживаются для просмотра. 
            Для открытия:</p>
            <ol style="margin: 10px 0 15px 20px; color: #856404; line-height: 1.8;">
                <li>Скачайте файл на ваш компьютер</li>
                <li>Откройте его в Microsoft Word или LibreOffice</li>
                <li>При необходимости сохраните в формате .docx</li>
            </ol>
            """
        elif recovery_method == "library_missing":
            recovery_text = """
            <h4>📚 Библиотека не установлена</h4>
            <p>Для просмотра содержимого .docx файлов требуется библиотека python-docx. 
            Вы можете:</p>
            <ul style="margin: 10px 0 15px 20px; color: #856404; line-height: 1.8;">
                <li>Скачать файл и открыть в Word на вашем компьютере</li>
                <li>Или установить библиотеку на сервер: <code>pip install python-docx</code></li>
            </ul>
            """
        else:
            recovery_text = ""

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Просмотр Word - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{ 
                    background: #f5f5f5; 
                    font-family: 'Segoe UI', Arial, sans-serif; 
                    display: flex; 
                    flex-direction: column; 
                    height: 100vh; 
                }}
                .toolbar {{
                    background: #2b579a;
                    color: white;
                    padding: 15px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.3);
                    flex-wrap: wrap;
                    gap: 10px;
                }}
                .toolbar h3 {{ margin: 0; font-size: 18px; }}
                .btn {{
                    padding: 10px 20px;
                    background: white;
                    color: #2b579a;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                    font-weight: bold;
                    transition: all 0.3s;
                    display: inline-flex;
                    align-items: center;
                    gap: 8px;
                }}
                .btn:hover {{ background: #e0e0e0; transform: translateY(-1px); }}
                .btn-download {{
                    background: #ffc107;
                    color: #333;
                    font-size: 16px;
                    padding: 12px 25px;
                }}
                .btn-download:hover {{
                    background: #ffb300;
                    box-shadow: 0 4px 15px rgba(255, 193, 7, 0.4);
                }}
                .content {{
                    flex: 1;
                    overflow-y: auto;
                    padding: 30px 40px;
                    background: white;
                    margin: 20px;
                    border-radius: 10px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                }}
                pre {{
                    white-space: pre-wrap;
                    word-wrap: break-word;
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 16px;
                    line-height: 1.8;
                    color: #333;
                }}
                .info {{
                    color: #666;
                    font-size: 14px;
                    margin-bottom: 20px;
                    padding-bottom: 15px;
                    border-bottom: 2px solid #2b579a;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    flex-wrap: wrap;
                    gap: 10px;
                }}
                .recovery-box {{
                    background: #fff3cd;
                    border: 2px solid #ffc107;
                    border-radius: 10px;
                    padding: 25px;
                    margin-bottom: 25px;
                }}
                .recovery-box h4 {{
                    color: #856404;
                    margin-bottom: 12px;
                    font-size: 18px;
                }}
                .recovery-box p {{
                    color: #856404;
                    margin-bottom: 15px;
                    font-size: 15px;
                    line-height: 1.7;
                }}
                .recovery-box ol, .recovery-box ul {{
                    color: #856404;
                    margin-bottom: 20px;
                    font-size: 15px;
                }}
                .recovery-box code {{
                    background: #ffeaa7;
                    padding: 2px 6px;
                    border-radius: 3px;
                    font-size: 13px;
                }}
                .download-section {{
                    text-align: center;
                    padding: 30px;
                }}
                .download-icon {{
                    font-size: 64px;
                    margin-bottom: 15px;
                }}
                .error-details {{
                    background: #ffeaa7;
                    padding: 12px 15px;
                    border-radius: 5px;
                    margin-top: 15px;
                    font-size: 13px;
                    color: #856404;
                    font-family: monospace;
                    white-space: pre-wrap;
                    word-break: break-all;
                }}
                ::-webkit-scrollbar {{
                    width: 8px;
                }}
                ::-webkit-scrollbar-track {{
                    background: #f1f1f1;
                    border-radius: 10px;
                }}
                ::-webkit-scrollbar-thumb {{
                    background: #2b579a;
                    border-radius: 10px;
                }}
                .status-badge {{
                    display: inline-block;
                    padding: 4px 12px;
                    border-radius: 12px;
                    font-size: 13px;
                    font-weight: bold;
                }}
                .status-ok {{
                    background: #d4edda;
                    color: #155724;
                }}
                .status-recovered {{
                    background: #fff3cd;
                    color: #856404;
                }}
                .status-error {{
                    background: #f8d7da;
                    color: #721c24;
                }}
                .text-section {{
                    margin-top: 25px;
                }}
                .text-section h3 {{
                    color: #2b579a;
                    margin-bottom: 15px;
                    font-size: 18px;
                }}
                .footer-tip {{
                    margin-top: 30px;
                    padding-top: 20px;
                    border-top: 2px solid #eee;
                    color: #999;
                    font-size: 12px;
                    display: flex;
                    justify-content: space-between;
                    flex-wrap: wrap;
                    gap: 10px;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <h3>📘 {filename}</h3>
                <div style="display: flex; gap: 10px; flex-wrap: wrap;">
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="content">
                <div class="info">
                    <div>
                        📄 Файл: <strong>{filename}</strong> | 
                        📏 Размер: <strong>{file_size} байт</strong>
                    </div>
                    {status_badge}
                </div>

                {f'''
                <div class="recovery-box">
                    {recovery_text}
                    <div class="download-section">
                        <div class="download-icon">📥</div>
                        <a href="/documents/{filename}" download class="btn btn-download">
                            💾 Скачать файл на компьютер
                        </a>
                        <p style="margin-top: 15px; font-size: 13px; color: #856404;">
                            После скачивания откройте файл в Microsoft Word для автоматического восстановления
                        </p>
                    </div>
                    {f'<div class="error-details">📋 Техническая информация: {error_msg}</div>' if error_msg else ''}
                </div>
                ''' if recovery_needed else ''}

                <div class="text-section">
                    <h3>{status_title}</h3>
                    <pre>{escaped_text}</pre>
                </div>

                <div class="footer-tip">
                    <span>💡 <strong>Совет:</strong> Microsoft Word автоматически восстанавливает поврежденные файлы при открытии</span>
                    <span>⌨️ <strong>Ctrl+S</strong> — скачать файл</span>
                </div>
            </div>

            <script>
                // Горячие клавиши
                document.addEventListener('keydown', function(e) {{
                    // Ctrl+S — скачать
                    if (e.ctrlKey && e.key === 's') {{
                        e.preventDefault();
                        window.location.href = '/documents/{filename}';
                    }}
                }});

                // Автофокус на кнопку скачивания если файл поврежден
                window.addEventListener('load', function() {{
                    const downloadBtn = document.querySelector('.btn-download');
                    if (downloadBtn) {{
                        // Подсвечиваем кнопку
                        downloadBtn.style.transform = 'scale(1.05)';
                        setTimeout(function() {{
                            downloadBtn.style.transform = 'scale(1)';
                        }}, 300);
                    }}
                }});
            </script>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ EXCEL ФАЙЛОВ (xlsx, xls) - ТОЛЬКО ЛОКАЛЬНО
    # ============================================
    elif file_extension in ['xlsx', 'xls']:
        # Пробуем прочитать данные из Excel
        table_html = ""
        try:
            if file_extension == 'xlsx':
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                    ws = wb.active

                    # Строим HTML таблицу (первые 100 строк и 20 колонок)
                    table_html = '<table class="excel-table"><thead><tr><th>#</th>'
                    max_col = min(ws.max_column or 20, 20)
                    max_row = min(ws.max_row or 100, 100)

                    # Заголовки
                    for col in range(1, max_col + 1):
                        cell_value = ws.cell(1, col).value
                        table_html += f'<th>{cell_value if cell_value is not None else ""}</th>'
                    table_html += '</tr></thead><tbody>'

                    # Данные
                    for row in range(2, max_row + 1):
                        table_html += '<tr>'
                        table_html += f'<td class="row-num">{row - 1}</td>'
                        for col in range(1, max_col + 1):
                            cell_value = ws.cell(row, col).value
                            table_html += f'<td>{cell_value if cell_value is not None else ""}</td>'
                        table_html += '</tr>'

                    table_html += '</tbody></table>'

                    if ws.max_row > 100 or ws.max_column > 20:
                        table_html += f'<p class="warning">⚠️ Показаны первые {max_row} строк и {max_col} колонок. Всего: {ws.max_row} строк, {ws.max_column} колонок.</p>'

                    wb.close()
                except ImportError:
                    table_html = '<p class="error">❌ Библиотека openpyxl не установлена. Установите: pip install openpyxl</p>'
                except Exception as e:
                    table_html = f'<p class="error">❌ Не удалось прочитать Excel файл: {str(e)}</p>'
            else:
                # Для .xls (старый формат)
                try:
                    import xlrd
                    wb = xlrd.open_workbook(file_path)
                    ws = wb.sheet_by_index(0)

                    table_html = '<table class="excel-table"><thead><tr><th>#</th>'
                    max_col = min(ws.ncols, 20)
                    max_row = min(ws.nrows, 100)

                    for col in range(max_col):
                        table_html += f'<th>{ws.cell_value(0, col)}</th>'
                    table_html += '</tr></thead><tbody>'

                    for row in range(1, max_row):
                        table_html += '<tr>'
                        table_html += f'<td class="row-num">{row}</td>'
                        for col in range(max_col):
                            table_html += f'<td>{ws.cell_value(row, col)}</td>'
                        table_html += '</tr>'

                    table_html += '</tbody></table>'

                    if ws.nrows > 100 or ws.ncols > 20:
                        table_html += f'<p class="warning">⚠️ Показаны первые {max_row} строк и {max_col} колонок. Всего: {ws.nrows} строк, {ws.ncols} колонок.</p>'

                except ImportError:
                    table_html = '<p class="error">❌ Библиотека xlrd не установлена. Установите: pip install xlrd</p>'
                except Exception as e:
                    table_html = f'<p class="error">❌ Не удалось прочитать файл .xls: {str(e)}</p>'
        except:
            table_html = '<p class="error">❌ Не удалось обработать файл.</p>'

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Просмотр Excel - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{ 
                    background: #f5f5f5; 
                    font-family: 'Segoe UI', Arial, sans-serif; 
                    display: flex; 
                    flex-direction: column; 
                    height: 100vh; 
                }}
                .toolbar {{
                    background: #217346;
                    color: white;
                    padding: 15px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.3);
                    flex-wrap: wrap;
                    gap: 10px;
                }}
                .toolbar h3 {{ margin: 0; font-size: 18px; }}
                .btn {{
                    padding: 8px 15px;
                    background: white;
                    color: #217346;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                    font-weight: bold;
                    transition: background 0.3s;
                }}
                .btn:hover {{ background: #e0e0e0; }}
                .content {{
                    flex: 1;
                    overflow: auto;
                    padding: 20px;
                }}
                .excel-table {{
                    border-collapse: collapse;
                    width: 100%;
                    background: white;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                    font-size: 13px;
                }}
                .excel-table th {{
                    background: #217346;
                    color: white;
                    padding: 10px;
                    border: 1px solid #1a5c37;
                    position: sticky;
                    top: 0;
                    z-index: 10;
                    white-space: nowrap;
                }}
                .excel-table td {{
                    padding: 8px 10px;
                    border: 1px solid #ddd;
                    white-space: nowrap;
                    max-width: 300px;
                    overflow: hidden;
                    text-overflow: ellipsis;
                }}
                .excel-table tr:hover td {{
                    background: #f0f7f3;
                }}
                .row-num {{
                    background: #f5f5f5;
                    color: #666;
                    text-align: center;
                    font-weight: bold;
                    width: 50px;
                }}
                .error {{
                    color: #e74c3c;
                    background: #ffeaa7;
                    padding: 15px;
                    border-radius: 8px;
                    margin: 20px;
                }}
                .warning {{
                    color: #856404;
                    background: #fff3cd;
                    padding: 10px 15px;
                    border-radius: 5px;
                    margin-top: 15px;
                }}
                .info {{
                    color: #666;
                    font-size: 14px;
                    margin-bottom: 20px;
                }}
                ::-webkit-scrollbar {{
                    width: 8px;
                    height: 8px;
                }}
                ::-webkit-scrollbar-track {{
                    background: #f1f1f1;
                }}
                ::-webkit-scrollbar-thumb {{
                    background: #217346;
                    border-radius: 10px;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <h3>📊 {filename}</h3>
                <div>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="content">
                <div class="info">
                    📊 Файл: <strong>{filename}</strong> | 
                    📏 Размер: <strong>{file_size} байт</strong> |
                    🔒 Локальный просмотр
                </div>
                {table_html}
            </div>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ POWERPOINT (pptx, ppt) - ТОЛЬКО ЛОКАЛЬНО
    # ============================================
    elif file_extension in ['pptx', 'ppt']:
        # Извлекаем текст из презентации
        slides_text = ""
        try:
            if file_extension == 'pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)

                    for i, slide in enumerate(prs.slides, 1):
                        slides_text += f'\n{"=" * 60}\n📄 Слайд {i}\n{"=" * 60}\n'
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slides_text += f'{shape.text}\n'
                            if shape.has_table:
                                table = shape.table
                                slides_text += '\n📊 Таблица:\n'
                                for row in table.rows:
                                    row_text = ' | '.join(cell.text for cell in row.cells)
                                    slides_text += f'  {row_text}\n'

                    if not slides_text:
                        slides_text = "(презентация не содержит текста)"
                except ImportError:
                    slides_text = "❌ Библиотека python-pptx не установлена.\nУстановите: pip install python-pptx"
                except Exception as e:
                    slides_text = f"❌ Не удалось прочитать презентацию: {str(e)}"
            else:
                slides_text = "⚠️ Файлы .ppt (старый формат) не поддерживаются.\nКонвертируйте в .pptx или используйте другую программу."
        except:
            slides_text = "❌ Не удалось извлечь текст из презентации."

        escaped_slides = html_module.escape(slides_text)

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Просмотр PowerPoint - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{ 
                    background: #f5f5f5; 
                    font-family: 'Segoe UI', Arial, sans-serif; 
                    display: flex; 
                    flex-direction: column; 
                    height: 100vh; 
                }}
                .toolbar {{
                    background: #d24726;
                    color: white;
                    padding: 15px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                    box-shadow: 0 2px 5px rgba(0,0,0,0.3);
                    flex-wrap: wrap;
                    gap: 10px;
                }}
                .toolbar h3 {{ margin: 0; font-size: 18px; }}
                .btn {{
                    padding: 8px 15px;
                    background: white;
                    color: #d24726;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                    font-weight: bold;
                    transition: background 0.3s;
                }}
                .btn:hover {{ background: #e0e0e0; }}
                .content {{
                    flex: 1;
                    overflow-y: auto;
                    padding: 30px 40px;
                    background: white;
                    margin: 20px;
                    border-radius: 10px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                }}
                pre {{
                    white-space: pre-wrap;
                    word-wrap: break-word;
                    font-family: 'Segoe UI', Arial, sans-serif;
                    font-size: 15px;
                    line-height: 1.7;
                    color: #333;
                }}
                .info {{
                    color: #666;
                    font-size: 14px;
                    margin-bottom: 20px;
                    padding-bottom: 15px;
                    border-bottom: 2px solid #d24726;
                }}
                ::-webkit-scrollbar {{
                    width: 8px;
                }}
                ::-webkit-scrollbar-track {{
                    background: #f1f1f1;
                }}
                ::-webkit-scrollbar-thumb {{
                    background: #d24726;
                    border-radius: 10px;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <h3>📽️ {filename}</h3>
                <div>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="content">
                <div class="info">
                    📽️ Файл: <strong>{filename}</strong> | 
                    📏 Размер: <strong>{file_size} байт</strong> |
                    📝 Локальный просмотр текста слайдов
                </div>
                <pre>{escaped_slides}</pre>
            </div>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ ВИДЕО ФАЙЛОВ (mp4, webm, avi, mov)
    # ============================================
    elif file_extension in ['mp4', 'webm', 'avi', 'mov', 'mkv']:
        mime_types = {
            'mp4': 'video/mp4',
            'webm': 'video/webm',
            'avi': 'video/x-msvideo',
            'mov': 'video/quicktime',
            'mkv': 'video/x-matroska'
        }
        mime_type = mime_types.get(file_extension, 'video/mp4')

        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Видео плеер - {filename}</title>
            <style>
                * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                body {{
                    background: #000;
                    font-family: Arial, sans-serif;
                    display: flex;
                    flex-direction: column;
                    height: 100vh;
                }}
                .toolbar {{
                    background: #1a1a1a;
                    color: white;
                    padding: 10px 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                }}
                .btn {{
                    padding: 8px 15px;
                    background: #667eea;
                    color: white;
                    border: none;
                    border-radius: 5px;
                    cursor: pointer;
                    text-decoration: none;
                    font-size: 14px;
                }}
                .btn:hover {{ background: #764ba2; }}
                .video-container {{
                    flex: 1;
                    display: flex;
                    justify-content: center;
                    align-items: center;
                }}
                video {{
                    max-width: 100%;
                    max-height: calc(100vh - 50px);
                }}
                .error-message {{
                    display: none;
                    text-align: center;
                    color: white;
                    padding: 50px;
                }}
            </style>
        </head>
        <body>
            <div class="toolbar">
                <h3>🎬 {filename}</h3>
                <div>
                    <a href="/documents/{filename}" download class="btn">💾 Скачать</a>
                    <button onclick="window.close()" class="btn">✖ Закрыть</button>
                </div>
            </div>
            <div class="video-container">
                <video controls autoplay id="videoPlayer">
                    <source src="/documents/{filename}" type="{mime_type}">
                    Ваш браузер не поддерживает видео элемент.
                </video>
            </div>
            <div class="error-message" id="errorMessage">
                <h2>❌ Ошибка воспроизведения</h2>
                <p>Файл может быть поврежден или формат не поддерживается.</p>
                <a href="/documents/{filename}" download class="btn">💾 Скачать файл</a>
            </div>
            <script>
                var video = document.getElementById('videoPlayer');
                video.addEventListener('error', function() {{
                    document.querySelector('.video-container').style.display = 'none';
                    document.getElementById('errorMessage').style.display = 'block';
                }});
            </script>
        </body>
        </html>
        '''
        return make_response(html_content)

    # ============================================
    # ДЛЯ ВСЕХ ОСТАЛЬНЫХ ТИПОВ ФАЙЛОВ
    # ============================================
    else:
        html_content = f'''
        <!DOCTYPE html>
        <html lang="ru">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Файл - {filename}</title>
            <style>
                body {{
                    font-family: Arial, sans-serif;
                    display: flex;
                    justify-content: center;
                    align-items: center;
                    height: 100vh;
                    margin: 0;
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                }}
                .container {{
                    text-align: center;
                    background: white;
                    padding: 40px;
                    border-radius: 10px;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.3);
                    max-width: 500px;
                }}
                .file-icon {{
                    font-size: 64px;
                    margin-bottom: 20px;
                }}
                h1 {{
                    color: #333;
                    margin-bottom: 15px;
                    word-wrap: break-word;
                }}
                p {{
                    color: #666;
                    margin-bottom: 20px;
                }}
                .btn {{
                    display: inline-block;
                    padding: 10px 20px;
                    background: #667eea;
                    color: white;
                    text-decoration: none;
                    border-radius: 5px;
                    margin: 5px;
                    transition: background 0.3s;
                    border: none;
                    cursor: pointer;
                    font-size: 14px;
                }}
                .btn:hover {{
                    background: #764ba2;
                }}
            </style>
        </head>
        <body>
            <div class="container">
                <div class="file-icon">📎</div>
                <h1>{filename}</h1>
                <p>Предпросмотр для файлов типа <strong>.{file_extension}</strong> не поддерживается.<br>
                Вы можете скачать файл для просмотра в соответствующей программе.</p>
                <a href="/documents/{filename}" download class="btn">💾 Скачать файл</a>
                <button onclick="window.close()" class="btn">✖ Закрыть</button>
            </div>
        </body>
        </html>
        '''
        return make_response(html_content)


@app.route('/documents/<path:filename>')
@login_required
def download_document(filename):
    """Скачивание документа"""
    import os
    documents_folder = os.path.join(app.root_path, 'documents')
    try:
        return send_from_directory(documents_folder, filename)
    except FileNotFoundError:
        abort(404, description="Файл не найден")


@app.route('/pcm-temp/<filename>')
@login_required
def serve_pcm_temp(filename):
    """Раздача конвертированного PCM -> WAV"""
    import tempfile
    import os

    temp_dir = tempfile.gettempdir()
    file_path = os.path.join(temp_dir, filename)

    if os.path.exists(file_path):
        return send_from_directory(temp_dir, filename)
    else:
        abort(404)


if __name__ == '__main__':
    app.run(debug=True, port=5000, host='0.0.0.0')

# from flask import Flask, render_template, jsonify, send_from_directory
#
# app = Flask(__name__)
#
# # Основные данные (в будущем замени на базу данных)
# coordinates = [
#     {"id": 1, "uid": 12, "latitude": 55.7558, "longitude": 37.6173, "name": "Москва"},
#     {"id": 2, "uid": 13,"latitude": 59.9343, "longitude": 30.3351, "name": "Санкт-Петербург"},
#     {"id": 3, "uid": 12,"latitude": 48.8566, "longitude": 2.3522, "name": "Париж"},
#     {"id": 4, "uid": 17,"latitude": 51.5074, "longitude": -0.1278, "name": "Лондон"},
#     {"id": 5, "uid": 12,"latitude": -40.7128, "longitude": -74.0060, "name": "Нью-Йорк"},
#     {"id": 6, "uid": 122,"latitude": 35.6762, "longitude": 139.6503, "name": "Токио"},
#     {"id": 7, "uid": 12,"latitude": -33.8688, "longitude": 151.2093, "name": "Сидней"},
#     {"id": 8, "uid": 11,"latitude": -23.5505, "longitude": -46.6333, "name": "Сан-Паулу"},
#     {"id": 9, "uid": 32,"latitude": 1.0760, "longitude": 72.8777, "name": "Мумбаи"},
#     {"id": 10, "uid": 2,"latitude": 0.0444, "longitude": 71.2357, "name": "Каир"},
#     {"id": 11, "uid": 122, "latitude": 50.04, "longitude": 30.2357, "name": "Кив"},
# ]
#
# @app.route('/')
# def index():
#     return render_template('index.html')
#
# @app.route('/api/coordinates')
# def get_coordinates():
#     return jsonify(coordinates)
#
# # Для корректной отдачи geojson
# @app.route('/static/data/<path:filename>')
# def serve_geojson(filename):
#     return send_from_directory('static/data', filename)
#
# if __name__ == '__main__':
#     app.run(debug=True, port=5000)

# from flask import Flask, render_template, jsonify, send_from_directory
# import pyodbc
#
# app = Flask(__name__)

#####
####
# # Конфигурация базы данных
# DB_CONFIG = {
#     'driver': '{ODBC Driver 17 for SQL Server}',
#     'server': 'localhost',
#     'database': 'my_database',
#     'trusted_connection': 'yes'
# }

# from flask import Flask, render_template, jsonify, send_from_directory
# import pyodbc
# from config import DB_CONFIG  # импортируем конфигурацию
#
# app = Flask(__name__)
#
#
# def get_coordinates_from_db():
#     try:
#         # Формируем строку подключения из импортированного конфига
#         if 'trusted_connection' in DB_CONFIG:
#             conn_str = (f"DRIVER={DB_CONFIG['driver']};"
#                         f"SERVER={DB_CONFIG['server']};"
#                         f"DATABASE={DB_CONFIG['database']};"
#                         f"Trusted_Connection=yes;")
#         else:
#             conn_str = (f"DRIVER={DB_CONFIG['driver']};"
#                         f"SERVER={DB_CONFIG['server']};"
#                         f"DATABASE={DB_CONFIG['database']};"
#                         f"UID={DB_CONFIG['uid']};"
#                         f"PWD={DB_CONFIG['pwd']};")
#
#         conn = pyodbc.connect(conn_str)
#         cursor = conn.cursor()
#         cursor.execute("SELECT id, latitude, longitude, name FROM your_table ORDER BY id")
#
#         coordinates = [{"id": row[0], "latitude": float(row[1]),
#                         "longitude": float(row[2]), "name": row[3]}
#                        for row in cursor.fetchall()]
#
#         conn.close()
#         return coordinates
#
#     except Exception as e:
#         print(f"Ошибка БД: {e}")
#         return []

#######

# Или для SQL-аутентификации:
# DB_CONFIG = {
#     'driver': '{ODBC Driver 17 for SQL Server}',
#     'server': 'localhost',
#     'database': 'my_database',
#     'uid': 'sa',
#     'pwd': 'your_password'
# }
# Настройки подключения к SQL Server
# DB_CONFIG = {
#     'driver': '{ODBC Driver 17 for SQL Server}',
#     'server': 'your_server',
#     'database': 'your_db',
#     'trusted_connection': 'yes'  # или используйте 'uid' и 'pwd'
# }
#
#
# def get_coordinates_from_db():
#     """Получение данных из SQL Server"""
#     try:
#         # Формируем строку подключения
#         if DB_CONFIG.get('trusted_connection'):
#             conn_str = (
#                 f"DRIVER={DB_CONFIG['driver']};"
#                 f"SERVER={DB_CONFIG['server']};"
#                 f"DATABASE={DB_CONFIG['database']};"
#                 f"Trusted_Connection=yes;"
#             )
#         else:
#             conn_str = (
#                 f"DRIVER={DB_CONFIG['driver']};"
#                 f"SERVER={DB_CONFIG['server']};"
#                 f"DATABASE={DB_CONFIG['database']};"
#                 f"UID={DB_CONFIG['uid']};"
#                 f"PWD={DB_CONFIG['pwd']};"
#             )
#
#         conn = pyodbc.connect(conn_str)
#         cursor = conn.cursor()
#         cursor.execute("SELECT id, latitude, longitude, name FROM your_table ORDER BY id")
#
#         coordinates = [
#             {"id": row[0], "latitude": float(row[1]), "longitude": float(row[2]), "name": row[3]}
#             for row in cursor.fetchall()
#         ]
#
#         conn.close()
#         return coordinates
#
#     except Exception as e:
#         print(f"Ошибка БД: {e}")
#         # Возвращаем тестовые данные в случае ошибки
#         return [
#             {"id": 1, "latitude": 55.7558, "longitude": 37.6173, "name": "Москва"},
#             {"id": 2, "latitude": 59.9343, "longitude": 30.3351, "name": "Санкт-Петербург"},
#         ]
#
#
# @app.route('/')
# def index():
#     return render_template('index.html')
#
#
# @app.route('/api/coordinates')
# def get_coordinates():
#     """API endpoint для получения координат из БД"""
#     coordinates = get_coordinates_from_db()
#     return jsonify(coordinates)
#
#
# @app.route('/static/data/<path:filename>')
# def serve_geojson(filename):
#     return send_from_directory('static/data', filename)
#
#
# if __name__ == '__main__':
#     app.run(debug=True, port=5000)